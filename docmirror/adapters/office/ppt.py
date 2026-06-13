# Copyright (c) 2026 ValueMap Global and contributors. All rights reserved.
# Author: Adam Lin <adamlin@valuemapglobal.com>
#
# This source code is licensed under the Apache 2.0 license found in the
# LICENSE file in the root directory of this source tree.

"""
PPT Adapter — .pptx → ParseResult
==================================

Extracts slide content via python-pptx. Legacy ``.ppt`` is transcoded upstream
by ``TranscodingGate`` (FCR); this adapter only receives ``.pptx``.
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import List

from docmirror.framework.base import BaseParser

logger = logging.getLogger(__name__)


class PPTAdapter(BaseParser):
    """PowerPoint (.pptx) adapter — python-pptx native parsing only."""

    async def to_parse_result(self, file_path: Path, **kwargs) -> ParseResult:
        """
        Parse a .pptx file into a ParseResult.

        Each slide maps to a PageContent. Slide titles become H2 TextBlocks,
        text shapes become BODY TextBlocks, and table shapes become TableBlocks.
        """
        from docmirror.models.errors import build_failure_result

        path = Path(file_path)
        if path.suffix.lower() != ".pptx":
            return build_failure_result(
                "FORMAT_REQUIRES_CONVERTER",
                f"PPTAdapter requires .pptx input; got {path.suffix}. "
                "Legacy .ppt is transcoded by the extraction pipeline.",
                file_path=str(path),
                file_type="ppt",
            )

        from pptx import Presentation

        from docmirror.models.entities.parse_result import (
            CellValue,
            DataType,
            PageContent,
            ParseResult,
            ParserInfo,
            TableBlock,
            TableRow,
            TextBlock,
            TextLevel,
        )

        logger.info(f"[PPTAdapter] Starting native extraction for presentation: {file_path}")
        prs = Presentation(str(file_path))

        pages: list[PageContent] = []

        for i, slide in enumerate(prs.slides):
            texts: list[TextBlock] = []
            tables: list[TableBlock] = []

            # Extract slide title (if present)
            if slide.shapes.title and slide.shapes.title.text:
                texts.append(
                    TextBlock(
                        content=slide.shapes.title.text,
                        level=TextLevel.H2,
                    )
                )

            # Extract content from all other shapes (skip the title shape)
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue

                # Text content from text frames
                if hasattr(shape, "text") and shape.text:
                    texts.append(TextBlock(content=shape.text, level=TextLevel.BODY))

                # Table content — typed CellValue
                if shape.has_table:
                    first_row = True
                    headers: list[str] = []
                    data_rows: list[TableRow] = []

                    for row in shape.table.rows:
                        cells = [CellValue(text=cell.text.strip(), data_type=DataType.TEXT) for cell in row.cells]
                        if first_row:
                            headers = [c.text for c in cells]
                            first_row = False
                        else:
                            if any(c.text for c in cells):
                                data_rows.append(TableRow(cells=cells, source_page=i))

                    if headers or data_rows:
                        tables.append(
                            TableBlock(
                                table_id=f"slide{i}_table{len(tables)}",
                                headers=headers,
                                rows=data_rows,
                                page=i,
                            )
                        )

            pages.append(PageContent(page_number=i, texts=texts, tables=tables))

        return ParseResult(
            pages=pages,
            parser_info=ParserInfo(
                parser_name="PPTAdapter",
                page_count=len(prs.slides),
                overall_confidence=1.0,
            ),
        )
